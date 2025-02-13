import os
import sys
import re
import uuid
import base64
from io import BytesIO, StringIO
from flask import Flask, render_template, request, send_file, redirect, url_for
from email import policy, message_from_file
from extract_msg import Message  # For .msg files
from weasyprint import HTML, CSS
from PyPDF2 import PdfMerger
import pypandoc  # For converting attachments to PDF
import pandas as pd
from pptx import Presentation  # For PowerPoint files
from openpyxl import load_workbook  # For Excel files

app = Flask(__name__)

def sanitize_html(html_content):
    """Sanitize HTML content to ensure it fits within the PDF page."""
    css = """
    <style>
        body { 
            font-family: Arial, sans-serif; 
            font-size: 12px; 
            margin: 0; 
            padding: 10px; 
            width: 100%; 
            overflow: hidden; 
        }
        pre, p { 
            white-space: pre-wrap; 
            word-wrap: break-word; 
            overflow-wrap: break-word; 
            max-width: 100%; 
        }
        img { 
            max-width: 100% !important; 
            height: auto !important; 
            display: block; 
            margin: 0 auto; 
        }
        table { 
            width: 100% !important; 
            max-width: 100% !important; 
            border-collapse: collapse; 
        }
        td, th { 
            border: 1px solid #ddd; 
            padding: 8px; 
            word-wrap: break-word; 
            max-width: 100%; 
        }
        .email-header { 
            margin-bottom: 20px; 
        }
        .email-header p { 
            margin: 5px 0; 
        }
        * { 
            box-sizing: border-box; 
            max-width: 100% !important; 
        }
    </style>
    """
    html_content = re.sub(r"<head>", f"<head>{css}", html_content, flags=re.IGNORECASE)
    return html_content

def convert_attachment_to_pdf(attachment_data, attachment_name):
    """Convert an attachment to PDF using pypandoc, pandas, or python-pptx."""
    try:
        # Save the attachment to a temporary file
        temp_input = f"/tmp/{uuid.uuid4()}_{attachment_name}"
        temp_output = f"/tmp/{uuid.uuid4()}.pdf"
        
        with open(temp_input, "wb") as f:
            f.write(attachment_data)
        
        # Handle Excel files
        if attachment_name.endswith(".xlsx"):
            # Read the Excel file into a DataFrame
            wb = load_workbook(temp_input)
            html_content = "<html><body>"
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                html_content += f"<h2>{sheet_name}</h2><table>"
                for row in sheet.iter_rows(values_only=True):
                    html_content += "<tr>"
                    for cell in row:
                        html_content += f"<td>{cell}</td>"
                    html_content += "</tr>"
                html_content += "</table>"
            html_content += "</body></html>"
            
            # Convert the HTML to PDF using WeasyPrint
            HTML(string=html_content).write_pdf(temp_output)
        
        # Handle PowerPoint files
        elif attachment_name.endswith(".pptx"):
            # Read the PowerPoint file
            prs = Presentation(temp_input)
            html_content = "<html><body>"
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        html_content += f"<p>{shape.text}</p>"
            html_content += "</body></html>"
            
            # Convert the HTML to PDF using WeasyPrint
            HTML(string=html_content).write_pdf(temp_output)
        
        # Handle Word files
        elif attachment_name.endswith(".docx"):
            # Use pypandoc for Word files
            pypandoc.convert_file(
                temp_input,
                "pdf",
                outputfile=temp_output,
                format="docx",
                extra_args=["--pdf-engine=pdflatex"]  # Specify the PDF engine
            )
        else:
            # Skip unsupported formats
            print(f"Unsupported file format: {attachment_name}")
            return None
        
        # Read the converted PDF
        with open(temp_output, "rb") as f:
            pdf_data = f.read()
        
        return pdf_data
    except Exception as e:
        print(f"Error converting {attachment_name} to PDF: {e}")
        return None
    finally:
        # Clean up temporary files
        if os.path.exists(temp_input):
            os.remove(temp_input)
        if os.path.exists(temp_output):
            os.remove(temp_output)

def msg_to_pdf(msg_file):
    """Convert .msg file to PDF."""
    msg = Message(msg_file)
    
    # Extract email headers
    headers = f"""
    <div class="email-header">
        <p><strong>From:</strong> {msg.sender}</p>
        <p><strong>To:</strong> {msg.to}</p>
        <p><strong>Subject:</strong> {msg.subject}</p>
    </div>
    """
    
    # Extract HTML content if available
    html_content = None
    if msg.htmlBody:
        html_content = msg.htmlBody
    else:
        # Fall back to plain text if no HTML content is found
        html_content = f"""
        <html>
            <head>
                <title>Email</title>
            </head>
            <body>
                {headers}
                <pre>{msg.body}</pre>
            </body>
        </html>
        """
    
    # Ensure html_content is a string (decode if it's bytes)
    if isinstance(html_content, bytes):
        html_content = html_content.decode("utf-8", errors="replace")
    
    # Inject headers into the HTML content
    if "<body>" in html_content:
        html_content = html_content.replace("<body>", f"<body>{headers}")
    
    # Handle embedded images
    for attachment in msg.attachments:
        if attachment.type == "image":
            # Encode the image as base64
            image_data = base64.b64encode(attachment.data).decode("utf-8")
            # Replace the cid reference in the HTML with the base64-encoded data
            cid = attachment.cid.strip("<>")
            if cid:
                html_content = html_content.replace(
                    f'cid:{cid}',
                    f'data:{attachment.mimetype};base64,{image_data}'
                )
    
    # Sanitize and generate PDF
    html_content = sanitize_html(html_content)
    pdf_bytes = BytesIO()
    HTML(string=html_content).write_pdf(
        pdf_bytes,
        stylesheets=[CSS(string="@page { size: A3 landscape; margin: 1cm; }")]
    )
    pdf_bytes.seek(0)
    return pdf_bytes, msg.attachments

def eml_to_pdf(eml_file):
    """Convert .eml file to PDF."""
    # Decode the bytes into a string
    eml_content = eml_file.read().decode("utf-8", errors="replace")
    # Use StringIO to simulate a file-like object for message_from_file
    eml_string_io = StringIO(eml_content)
    msg = message_from_file(eml_string_io, policy=policy.default)
    
    # Extract email headers
    headers = f"""
    <div class="email-header">
        <p><strong>From:</strong> {msg['From']}</p>
        <p><strong>To:</strong> {msg['To']}</p>
        <p><strong>Subject:</strong> {msg['Subject']}</p>
    </div>
    """
    
    # Extract HTML content from the email
    html_content = None
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/html":
                html_content = part.get_payload(decode=True).decode("utf-8", errors="replace")
                break
    else:
        if msg.get_content_type() == "text/html":
            html_content = msg.get_payload(decode=True).decode("utf-8", errors="replace")
    
    # If no HTML content is found, fall back to plain text
    if not html_content:
        plain_text = msg.get_payload(decode=True).decode("utf-8", errors="replace")
        html_content = f"""
        <html>
            <head>
                <title>Email</title>
            </head>
            <body>
                {headers}
                <pre>{plain_text}</pre>
            </body>
        </html>
        """
    
    # Ensure html_content is a string (decode if it's bytes)
    if isinstance(html_content, bytes):
        html_content = html_content.decode("utf-8", errors="replace")
    
    # Inject headers into the HTML content
    if "<body>" in html_content:
        html_content = html_content.replace("<body>", f"<body>{headers}")
    
    # Handle inline images and attachments
    attachments = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_maintype() == "image":
                # Extract the image data
                image_data = part.get_payload(decode=True)
                # Encode the image as base64
                image_base64 = base64.b64encode(image_data).decode("utf-8")
                # Replace the cid reference in the HTML with the base64-encoded data
                cid = part.get("Content-ID", "").strip("<>")
                if cid:
                    html_content = html_content.replace(
                        f'cid:{cid}',
                        f'data:{part.get_content_type()};base64,{image_base64}'
                    )
            elif part.get_filename():
                # Extract attachments
                attachment_name = part.get_filename()
                attachment_data = part.get_payload(decode=True)
                attachment_type = part.get_content_type()
                attachments.append({
                    "name": attachment_name,
                    "data": attachment_data,
                    "type": attachment_type
                })
                print(f"Found attachment: {attachment_name} (MIME type: {attachment_type})")
    
    # Sanitize and generate PDF
    html_content = sanitize_html(html_content)
    pdf_bytes = BytesIO()
    HTML(string=html_content).write_pdf(
        pdf_bytes,
        stylesheets=[CSS(string="@page { size: A3 landscape; margin: 1cm; }")]
    )
    pdf_bytes.seek(0)
    return pdf_bytes, attachments

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        # Check if a file was uploaded
        if "file" not in request.files:
            return redirect(request.url)
        
        file = request.files["file"]
        if file.filename == "":
            return redirect(request.url)
        
        if file and (file.filename.endswith(".eml") or file.filename.endswith(".msg")):
            # Process the file in memory
            file_stream = BytesIO(file.read())
            if file.filename.endswith(".eml"):
                pdf_bytes, attachments = eml_to_pdf(file_stream)
            elif file.filename.endswith(".msg"):
                pdf_bytes, attachments = msg_to_pdf(file_stream)
            
            # Merge attachments into the PDF
            if attachments:
                merger = PdfMerger()
                merger.append(pdf_bytes)
                
                for attachment in attachments:
                    if attachment["type"] == "application/pdf":
                        # Directly append PDF attachments
                        merger.append(BytesIO(attachment["data"]))
                    else:
                        # Convert non-PDF attachments to PDF
                        print(f"Converting: {attachment['name']}")
                        pdf_data = convert_attachment_to_pdf(attachment["data"], attachment["name"])
                        if pdf_data:
                            merger.append(BytesIO(pdf_data))
                        else:
                            print(f"Skipping unsupported attachment: {attachment['name']}")
                   
                # Save the merged PDF to a BytesIO object
                merged_pdf = BytesIO()
                merger.write(merged_pdf)
                merged_pdf.seek(0)
                pdf_bytes = merged_pdf
            
            # Return the PDF as a downloadable file
            return send_file(
                pdf_bytes,
                as_attachment=True,
                download_name=f"{uuid.uuid4()}.pdf",
                mimetype="application/pdf"
            )
    
    return render_template("index.html")

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
